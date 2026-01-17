"""
Executive Summary Section - Enhanced with Metrics
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF

from .base import BaseSection
from .charts import ChartGenerator


class ExecutiveSummarySection(BaseSection):
    """
    Generates the Executive Summary slide from LLM summaries.

    Enhanced with:
    - Key metric cards at top
    - Condensed key insight
    - Visual hierarchy

    Requires: llm_summaries.json (overall_summary.executive_summary)
    """

    def get_slide_count(self) -> int:
        return 1

    def _get_summary_data(self) -> dict:
        """Extract executive summary and key metrics from LLM data."""
        llm_data = self.data.get('llm_summaries', {})
        overall = llm_data.get('overall_summary', {})
        metadata = llm_data.get('metadata', {})

        summary_text = overall.get('executive_summary', 'Executive summary not available. Run llm_batch_summarizer.py to generate.')

        # Get key metrics
        priorities = overall.get('strategic_priorities', [])
        question_summaries = llm_data.get('question_summaries', [])

        # Calculate sentiment breakdown
        sentiments = {'Positive': 0, 'Neutral': 0, 'Negative': 0}
        for q in question_summaries:
            s = q.get('sentiment_analysis', {}).get('overall', 'Neutral')
            if s in sentiments:
                sentiments[s] += 1

        total = sum(sentiments.values()) or 1
        positive_pct = round(sentiments['Positive'] / total * 100)

        return {
            'summary_text': summary_text,
            'total_responses': metadata.get('total_responses', 0),
            'total_questions': metadata.get('total_questions', 0),
            'priority_count': len(priorities),
            'positive_pct': positive_pct,
        }

    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        data = self._get_summary_data()
        charts = ChartGenerator()

        slide = self.add_content_slide(prs, "Executive Summary")

        # Generate metric cards at top
        metrics = [
            (f"{data['total_responses']:,}", "Responses", "primary"),
            (f"{data['positive_pct']}%", "Positive", "positive"),
            (str(data['priority_count']), "Priorities", "accent"),
            (str(data['total_questions']), "Questions", "secondary"),
        ]
        metrics_path = charts.metric_cards(metrics, figsize=(11, 1.8))
        self.add_chart_image(slide, metrics_path, 0.3, 1.2, 12.5, 1.5)

        # Key insight highlight box
        summary_text = data['summary_text']
        paragraphs = [p.strip() for p in summary_text.split('\n\n') if p.strip()]

        if paragraphs:
            # First paragraph as key insight (highlighted)
            key_insight = paragraphs[0][:300]
            if len(paragraphs[0]) > 300:
                key_insight += '...'

            # Highlight box
            highlight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.3), Inches(2.9), Inches(12.5), Inches(1.5)
            )
            highlight_box.fill.solid()
            highlight_box.fill.fore_color.rgb = self.colors['background']
            highlight_box.line.color.rgb = self.colors['secondary']

            self.add_text_box(
                slide, 0.5, 2.95, 0.3, 0.3,
                "KEY",
                font_size=8, bold=True, color='accent'
            )

            self.add_text_box(
                slide, 0.9, 3.0, 11.5, 1.3,
                key_insight,
                font_size=12, color='text'
            )

            # Additional context (remaining paragraphs)
            if len(paragraphs) > 1:
                additional_text = ' '.join(paragraphs[1:3])[:400]
                if len(' '.join(paragraphs[1:3])) > 400:
                    additional_text += '...'

                self.add_text_box(
                    slide, 0.5, 4.6, 12.0, 2.0,
                    additional_text,
                    font_size=11, color='text'
                )

        # Data source footer
        self.add_text_box(
            slide, 0.5, 6.8, 12.0, 0.4,
            "Analysis powered by Claude Opus 4.5 | Data from Mentimeter Survey",
            font_size=9, color='neutral', align='center'
        )

        charts.cleanup()
        return prs

    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        data = self._get_summary_data()

        self.pdf_add_content_page(pdf, "Executive Summary")

        # Metric row
        pdf.set_font('Helvetica', 'B', 20)
        metrics = [
            (f"{data['total_responses']:,}", "Responses", self.pdf_colors['primary']),
            (f"{data['positive_pct']}%", "Positive", self.pdf_colors['positive']),
            (str(data['priority_count']), "Priorities", self.pdf_colors['accent']),
        ]

        for value, label, color in metrics:
            pdf.set_text_color(*color)
            pdf.cell(60, 12, value, align='C')

        pdf.set_x(10)
        pdf.set_y(pdf.get_y() + 12)

        pdf.set_font('Helvetica', '', 9)
        pdf.set_text_color(*self.pdf_colors['neutral'])
        for _, label, _ in metrics:
            pdf.cell(60, 6, label, align='C')

        pdf.ln(10)

        # Summary text
        summary_text = data['summary_text']
        paragraphs = [p.strip() for p in summary_text.split('\n\n') if p.strip()]

        if paragraphs:
            # Key insight
            pdf.set_font('Helvetica', 'B', 10)
            pdf.set_text_color(*self.pdf_colors['accent'])
            pdf.cell(0, 8, "KEY INSIGHT:", new_x='LMARGIN', new_y='NEXT')

            pdf.set_font('Helvetica', '', 11)
            pdf.set_text_color(*self.pdf_colors['text'])
            self.pdf_add_paragraph(pdf, paragraphs[0])

            # Additional context
            for para in paragraphs[1:3]:
                self.pdf_add_paragraph(pdf, para, font_size=10)

        return pdf
