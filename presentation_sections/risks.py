"""
Critical Risks Section
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF

from .base import BaseSection


class RisksSection(BaseSection):
    """
    Generates the Critical Risks slide from LLM summaries.

    Requires: llm_summaries.json (overall_summary.critical_risks)

    Shows three risk categories:
    - People risks (retention, satisfaction)
    - Revenue risks (customer impact)
    - Competitive risks (market position)
    """

    def get_slide_count(self) -> int:
        return 1

    def _get_risks(self) -> dict:
        """Extract risks from LLM data."""
        llm_data = self.data.get('llm_summaries', {})
        overall = llm_data.get('overall_summary', {})
        risks = overall.get('critical_risks', {})

        if not risks:
            return {
                'people': 'Risk data not available. Run llm_batch_summarizer.py to generate.',
                'revenue': 'Risk data not available.',
                'competitive': 'Risk data not available.'
            }

        return risks

    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        risks = self._get_risks()

        slide = self.add_content_slide(prs, "Critical Risks")

        risk_items = [
            ('People', risks.get('people', 'N/A'), 'negative'),
            ('Revenue', risks.get('revenue', 'N/A'), 'accent'),
            ('Competitive', risks.get('competitive', 'N/A'), 'secondary'),
        ]

        # Three-column layout
        col_width = 4.0
        x_positions = [0.5, 4.7, 8.9]

        for i, (title, content, color) in enumerate(risk_items):
            x = x_positions[i]

            # Risk category header
            header_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(1.3), Inches(col_width), Inches(0.6)
            )
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = self.colors[color]
            header_box.line.fill.background()

            header_text = header_box.text_frame
            header_text.paragraphs[0].text = f"  {title} Risks"
            header_text.paragraphs[0].font.size = Pt(14)
            header_text.paragraphs[0].font.bold = True
            header_text.paragraphs[0].font.color.rgb = self.colors['white']

            # Risk content
            self.add_text_box(
                slide, x, 2.1, col_width, 4.5,
                content, font_size=11, color='text'
            )

        return prs

    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        risks = self._get_risks()

        self.pdf_add_content_page(pdf, "Critical Risks")

        risk_items = [
            ('People Risks', risks.get('people', 'N/A'), self.pdf_colors['negative']),
            ('Revenue Risks', risks.get('revenue', 'N/A'), self.pdf_colors['accent']),
            ('Competitive Risks', risks.get('competitive', 'N/A'), self.pdf_colors['secondary']),
        ]

        for title, content, color in risk_items:
            # Category header
            pdf.set_font('Helvetica', 'B', 12)
            pdf.set_fill_color(*color)
            pdf.set_text_color(*self.pdf_colors['white'])
            pdf.cell(0, 8, f"  {title}", fill=True, new_x='LMARGIN', new_y='NEXT')

            # Content
            pdf.set_text_color(*self.pdf_colors['text'])
            pdf.set_font('Helvetica', '', 10)
            pdf.multi_cell(0, 5, content)
            pdf.ln(6)

        return pdf
