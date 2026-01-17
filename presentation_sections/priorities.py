"""
Top 5 Priorities Section - Enhanced with Visual Ranking
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF

from .base import BaseSection
from .charts import ChartGenerator


class PrioritiesSection(BaseSection):
    """
    Generates the Top 5 Strategic Priorities slides from LLM summaries.

    Requires: llm_summaries.json (overall_summary.strategic_priorities)

    Creates 2 slides:
    - Slide 1: Priorities 1-3 with rationale
    - Slide 2: Priorities 4-5 with rationale
    """

    def get_slide_count(self) -> int:
        return 2

    def _get_priorities(self) -> list:
        """Extract priorities from LLM data."""
        llm_data = self.data.get('llm_summaries', {})
        overall = llm_data.get('overall_summary', {})
        priorities = overall.get('strategic_priorities', [])

        if not priorities:
            return [{
                'rank': i,
                'priority': f'Priority {i} not available',
                'rationale': 'Run llm_batch_summarizer.py to generate priorities.'
            } for i in range(1, 6)]

        return priorities

    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        priorities = self._get_priorities()
        charts = ChartGenerator()

        # Slide 1: Visual Priority Overview
        slide1 = self.add_content_slide(prs, "Top 5 Strategic Priorities")

        # Create priority ranking data for chart
        chart_data = []
        for p in priorities[:5]:
            rank = p.get('rank', 1)
            title = p.get('priority', 'Unknown')[:35]
            # Assign importance based on rank
            importance = 'HIGH' if rank <= 2 else ('MEDIUM' if rank <= 4 else 'LOW')
            chart_data.append((rank, title, importance))

        # Generate priority ranking chart
        if chart_data:
            chart_path = charts.priority_ranking(chart_data, figsize=(6, 5))
            self.add_chart_image(slide1, chart_path, 0.3, 1.3, 6.5, 5)

        # Add top 3 rationale summaries on right side
        y_pos = 1.4
        self.add_text_box(
            slide1, 7.0, 1.2, 5.5, 0.4,
            "Key Rationale", font_size=14, bold=True, color='secondary'
        )

        for p in priorities[:3]:
            rank = p.get('rank', '?')
            rationale = p.get('rationale', '')[:120]
            if len(p.get('rationale', '')) > 120:
                rationale += '...'

            self.add_text_box(
                slide1, 7.0, y_pos + 0.4, 5.5, 1.3,
                f"{rank}. {rationale}",
                font_size=10, color='text'
            )
            y_pos += 1.5

        # Slide 2: Priority Details (4-5 with full rationale)
        slide2 = self.add_content_slide(prs, "Strategic Priorities - Details")
        self._add_priority_content(slide2, priorities[3:5], start_y=1.3)

        charts.cleanup()
        return prs

    def _add_priority_content(self, slide, priorities: list, start_y: float):
        """Add priority items to a slide."""
        y_pos = start_y

        for priority in priorities:
            rank = priority.get('rank', '?')
            title = priority.get('priority', 'Unknown')
            rationale = priority.get('rationale', '')

            # Priority number badge
            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.5)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.colors['accent']
            badge.line.fill.background()

            # Badge text (rank number)
            badge_text = badge.text_frame
            badge_text.paragraphs[0].text = str(rank)
            badge_text.paragraphs[0].font.size = Pt(16)
            badge_text.paragraphs[0].font.bold = True
            badge_text.paragraphs[0].font.color.rgb = self.colors['white']
            from pptx.enum.text import PP_ALIGN
            badge_text.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Priority title
            self.add_text_box(
                slide, 1.2, y_pos - 0.05, 11.5, 0.5,
                title, font_size=16, bold=True, color='primary'
            )

            # Rationale
            self.add_text_box(
                slide, 1.2, y_pos + 0.5, 11.5, 1.2,
                rationale, font_size=12, color='text'
            )

            y_pos += 2.0

    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        priorities = self._get_priorities()

        # Page 1: Priorities 1-3
        self.pdf_add_content_page(pdf, "Top 5 Strategic Priorities (1-3)")
        self._add_pdf_priorities(pdf, priorities[:3])

        # Page 2: Priorities 4-5
        self.pdf_add_content_page(pdf, "Top 5 Strategic Priorities (4-5)")
        self._add_pdf_priorities(pdf, priorities[3:5])

        return pdf

    def _add_pdf_priorities(self, pdf: FPDF, priorities: list):
        """Add priority items to PDF."""
        for priority in priorities:
            rank = priority.get('rank', '?')
            title = self._sanitize_text(priority.get('priority', 'Unknown'))
            rationale = self._sanitize_text(priority.get('rationale', ''))

            # Priority header with number
            pdf.set_font('Helvetica', 'B', 14)
            pdf.set_text_color(*self.pdf_colors['accent'])
            pdf.cell(10, 8, f"{rank}.")
            pdf.set_text_color(*self.pdf_colors['primary'])
            pdf.multi_cell(0, 8, title)

            # Rationale
            pdf.set_font('Helvetica', '', 11)
            pdf.set_text_color(*self.pdf_colors['text'])
            pdf.set_x(15)
            pdf.multi_cell(0, 6, rationale)
            pdf.ln(6)
