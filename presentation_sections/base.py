"""
Base Section - Abstract base class for all presentation sections.
"""

from abc import ABC, abstractmethod
from typing import Dict, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF


# Shared color scheme
COLORS = {
    'primary': RGBColor(0x1E, 0x3A, 0x5F),      # Dark blue
    'secondary': RGBColor(0x3D, 0x5A, 0x80),    # Medium blue
    'accent': RGBColor(0xE0, 0x7A, 0x5F),       # Coral
    'positive': RGBColor(0x2E, 0x7D, 0x32),     # Green
    'negative': RGBColor(0xC6, 0x28, 0x28),     # Red
    'neutral': RGBColor(0x75, 0x75, 0x75),      # Gray
    'background': RGBColor(0xF5, 0xF5, 0xF5),   # Light gray
    'text': RGBColor(0x21, 0x21, 0x21),         # Dark gray
    'white': RGBColor(0xFF, 0xFF, 0xFF),        # White
}

# PDF color scheme (RGB tuples 0-255)
PDF_COLORS = {
    'primary': (30, 58, 95),
    'secondary': (61, 90, 128),
    'accent': (224, 122, 95),
    'positive': (46, 125, 50),
    'negative': (198, 40, 40),
    'neutral': (117, 117, 117),
    'background': (245, 245, 245),
    'text': (33, 33, 33),
    'white': (255, 255, 255),
}


class BaseSection(ABC):
    """
    Abstract base class for presentation sections.

    Each section must implement:
    - generate_pptx_slides(): Add slides to PowerPoint
    - generate_pdf_content(): Add pages to PDF
    - get_slide_count(): Return estimated slide count
    """

    def __init__(self, data_sources: Dict[str, Any], config: Optional[Dict] = None):
        """
        Initialize the section.

        Args:
            data_sources: Dict containing raw_data, llm_summaries, etc.
            config: Section-specific configuration
        """
        self.data = data_sources
        self.config = config or {}
        self.colors = COLORS
        self.pdf_colors = PDF_COLORS

    @abstractmethod
    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        """
        Add slides to the PowerPoint presentation.

        Args:
            prs: python-pptx Presentation object

        Returns:
            Modified Presentation object
        """
        pass

    @abstractmethod
    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        """
        Add pages to the PDF document.

        Args:
            pdf: FPDF object

        Returns:
            Modified FPDF object
        """
        pass

    @abstractmethod
    def get_slide_count(self) -> int:
        """Return the number of slides this section generates."""
        pass

    # Helper methods for PowerPoint

    def add_title_slide(self, prs: Presentation, title: str, subtitle: str = '') -> Presentation:
        """Add a title slide layout."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['primary']
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self.colors['secondary']
            subtitle_para.alignment = PP_ALIGN.CENTER

        return prs

    def add_content_slide(self, prs: Presentation, title: str) -> Any:
        """Add a content slide with title, return the slide object."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.0)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.colors['primary']
        header.line.fill.background()

        # Add title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['white']

        return slide

    def add_text_box(self, slide, left: float, top: float, width: float, height: float,
                     text: str, font_size: int = 14, bold: bool = False,
                     color: str = 'text', align: str = 'left') -> Any:
        """Add a text box to a slide."""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = text_box.text_frame
        frame.word_wrap = True

        para = frame.paragraphs[0]
        para.text = text
        para.font.size = Pt(font_size)
        para.font.bold = bold
        para.font.color.rgb = self.colors.get(color, self.colors['text'])

        if align == 'center':
            para.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            para.alignment = PP_ALIGN.RIGHT

        return text_box

    def add_bullet_list(self, slide, left: float, top: float, width: float, height: float,
                        items: list, font_size: int = 14) -> Any:
        """Add a bullet list to a slide."""
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        frame = text_box.text_frame
        frame.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                para = frame.paragraphs[0]
            else:
                para = frame.add_paragraph()

            para.text = f"  {item}"
            para.font.size = Pt(font_size)
            para.font.color.rgb = self.colors['text']
            para.level = 0

        return text_box

    def add_chart_image(self, slide, image_path: str, left: float, top: float,
                        width: float, height: float) -> Any:
        """
        Add a chart image to a slide.

        Args:
            slide: PowerPoint slide object
            image_path: Path to PNG image file
            left, top: Position in inches
            width, height: Size in inches

        Returns:
            Picture shape object
        """
        return slide.shapes.add_picture(
            image_path,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

    # Helper methods for PDF

    def pdf_add_title_page(self, pdf: FPDF, title: str, subtitle: str = '') -> FPDF:
        """Add a title page to PDF."""
        pdf.add_page()

        # Title
        pdf.set_font('Helvetica', 'B', 36)
        pdf.set_text_color(*self.pdf_colors['primary'])
        pdf.set_y(80)
        pdf.cell(0, 20, title, align='C', new_x='LMARGIN', new_y='NEXT')

        # Subtitle
        if subtitle:
            pdf.set_font('Helvetica', '', 18)
            pdf.set_text_color(*self.pdf_colors['secondary'])
            pdf.cell(0, 15, subtitle, align='C', new_x='LMARGIN', new_y='NEXT')

        return pdf

    def pdf_add_content_page(self, pdf: FPDF, title: str) -> FPDF:
        """Add a content page with header to PDF."""
        pdf.add_page()

        # Header bar
        pdf.set_fill_color(*self.pdf_colors['primary'])
        pdf.rect(0, 0, 297, 18, 'F')

        # Title
        pdf.set_font('Helvetica', 'B', 16)
        pdf.set_text_color(*self.pdf_colors['white'])
        pdf.set_xy(10, 4)
        pdf.cell(0, 10, title)

        # Reset for content
        pdf.set_text_color(*self.pdf_colors['text'])
        pdf.set_y(25)

        return pdf

    def _sanitize_text(self, text: str) -> str:
        """Sanitize text for PDF output (replace Unicode chars with ASCII equivalents)."""
        replacements = {
            '\u2014': '-',   # em-dash
            '\u2013': '-',   # en-dash
            '\u2018': "'",   # left single quote
            '\u2019': "'",   # right single quote
            '\u201c': '"',   # left double quote
            '\u201d': '"',   # right double quote
            '\u2026': '...', # ellipsis
            '\u2022': '*',   # bullet
            '\u00a0': ' ',   # non-breaking space
            '\u00b7': '*',   # middle dot
        }
        for unicode_char, ascii_char in replacements.items():
            text = text.replace(unicode_char, ascii_char)
        # Remove any remaining non-latin1 characters
        return text.encode('latin-1', errors='replace').decode('latin-1')

    def pdf_add_paragraph(self, pdf: FPDF, text: str, font_size: int = 11) -> FPDF:
        """Add a paragraph of text to PDF."""
        pdf.set_font('Helvetica', '', font_size)
        pdf.set_text_color(*self.pdf_colors['text'])
        pdf.multi_cell(0, 6, self._sanitize_text(text))
        pdf.ln(4)
        return pdf

    def pdf_add_bullet_list(self, pdf: FPDF, items: list, font_size: int = 11) -> FPDF:
        """Add a bullet list to PDF."""
        pdf.set_font('Helvetica', '', font_size)
        pdf.set_text_color(*self.pdf_colors['text'])

        for item in items:
            pdf.cell(10, 6, '*')  # Bullet character (ASCII safe)
            pdf.multi_cell(0, 6, self._sanitize_text(item))

        pdf.ln(4)
        return pdf
