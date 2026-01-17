"""
Title Slide Section
"""

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF

from .base import BaseSection


class TitleSlideSection(BaseSection):
    """
    Generates the presentation title slide.

    Config options:
        - title: Main presentation title (default: "Presales Survey Analysis")
        - subtitle: Subtitle text (default: "International All-Hands Insights")
        - presenter: Presenter name (optional)
        - date: Date string (default: today's date)
    """

    def get_slide_count(self) -> int:
        return 1

    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        title = self.config.get('title', 'Presales Survey Analysis')
        subtitle = self.config.get('subtitle', 'International All-Hands Insights')
        presenter = self.config.get('presenter', '')
        date = self.config.get('date', datetime.now().strftime('%B %Y'))

        # Add title slide
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background accent bar at bottom
        accent_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.colors['primary']
        accent_bar.line.fill.background()

        # Main title
        self.add_text_box(
            slide, 0.5, 2.0, 12.333, 1.5,
            title, font_size=48, bold=True, color='primary', align='center'
        )

        # Subtitle
        self.add_text_box(
            slide, 0.5, 3.8, 12.333, 0.8,
            subtitle, font_size=24, color='secondary', align='center'
        )

        # Date and presenter
        footer_text = date
        if presenter:
            footer_text = f"{presenter}  |  {date}"

        self.add_text_box(
            slide, 0.5, 6.7, 12.333, 0.5,
            footer_text, font_size=14, color='white', align='center'
        )

        return prs

    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        title = self.config.get('title', 'Presales Survey Analysis')
        subtitle = self.config.get('subtitle', 'International All-Hands Insights')
        presenter = self.config.get('presenter', '')
        date = self.config.get('date', datetime.now().strftime('%B %Y'))

        pdf.add_page()

        # Background accent at bottom
        pdf.set_fill_color(*self.pdf_colors['primary'])
        pdf.rect(0, 180, 297, 30, 'F')

        # Main title
        pdf.set_font('Helvetica', 'B', 36)
        pdf.set_text_color(*self.pdf_colors['primary'])
        pdf.set_y(70)
        pdf.cell(0, 20, title, align='C', new_x='LMARGIN', new_y='NEXT')

        # Subtitle
        pdf.set_font('Helvetica', '', 20)
        pdf.set_text_color(*self.pdf_colors['secondary'])
        pdf.cell(0, 15, subtitle, align='C', new_x='LMARGIN', new_y='NEXT')

        # Footer
        footer_text = date
        if presenter:
            footer_text = f"{presenter}  |  {date}"

        pdf.set_font('Helvetica', '', 12)
        pdf.set_text_color(*self.pdf_colors['white'])
        pdf.set_y(188)
        pdf.cell(0, 10, footer_text, align='C')

        return pdf
