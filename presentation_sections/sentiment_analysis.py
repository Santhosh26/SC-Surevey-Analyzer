"""
Sentiment Analysis Section - Enhanced with Charts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from fpdf import FPDF

from .base import BaseSection
from .charts import ChartGenerator


class SentimentAnalysisSection(BaseSection):
    """
    Generates Sentiment Analysis slides showing overall sentiment breakdown.

    Enhanced with:
    - Donut chart for overall sentiment
    - Per-question sentiment indicators

    Requires: llm_summaries.json (for question-level sentiment)
    """

    def get_slide_count(self) -> int:
        return 1

    def _get_sentiment_data(self) -> dict:
        """Get sentiment breakdown from LLM summaries."""
        llm_data = self.data.get('llm_summaries', {})

        # Try to aggregate from question summaries
        summaries = llm_data.get('question_summaries', [])
        question_sentiments = []

        if summaries:
            sentiments = {'Positive': 0, 'Neutral': 0, 'Negative': 0}
            for s in summaries:
                sentiment = s.get('sentiment_analysis', {}).get('overall', 'Neutral')
                if sentiment in sentiments:
                    sentiments[sentiment] += 1

                # Store per-question sentiment
                q_short = s.get('question', '')[:25]
                if len(s.get('question', '')) > 25:
                    q_short += '...'
                question_sentiments.append((q_short, sentiment))

            total = sum(sentiments.values()) or 1
            return {
                'positive_pct': round(sentiments['Positive'] / total * 100),
                'neutral_pct': round(sentiments['Neutral'] / total * 100),
                'negative_pct': round(sentiments['Negative'] / total * 100),
                'positive_count': sentiments['Positive'],
                'neutral_count': sentiments['Neutral'],
                'negative_count': sentiments['Negative'],
                'question_count': total,
                'question_sentiments': question_sentiments,
            }

        return {
            'positive_pct': 33,
            'neutral_pct': 34,
            'negative_pct': 33,
            'positive_count': 0,
            'neutral_count': 0,
            'negative_count': 0,
            'question_count': 0,
            'question_sentiments': [],
        }

    def generate_pptx_slides(self, prs: Presentation) -> Presentation:
        data = self._get_sentiment_data()
        charts = ChartGenerator()

        slide = self.add_content_slide(prs, "Sentiment Analysis")

        # Generate donut chart (left side)
        donut_path = charts.donut_chart(
            values=[data['positive_count'], data['neutral_count'], data['negative_count']],
            labels=['Positive', 'Neutral', 'Negative'],
            center_text=f"{data['question_count']}\nQuestions",
            figsize=(4.5, 4.5)
        )
        self.add_chart_image(slide, donut_path, 0.3, 1.5, 5, 5)

        # Key insight text (right side)
        self.add_text_box(
            slide, 5.8, 1.5, 7.0, 0.5,
            "Overall Sentiment Distribution",
            font_size=16, bold=True, color='primary'
        )

        # Sentiment breakdown cards
        y_pos = 2.2
        sentiments = [
            ('Positive', data['positive_pct'], data['positive_count'], 'positive'),
            ('Neutral', data['neutral_pct'], data['neutral_count'], 'neutral'),
            ('Negative', data['negative_pct'], data['negative_count'], 'negative'),
        ]

        for label, pct, count, color in sentiments:
            # Color indicator
            indicator = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(5.8), Inches(y_pos), Inches(0.3), Inches(0.5)
            )
            indicator.fill.solid()
            indicator.fill.fore_color.rgb = self.colors[color]
            indicator.line.fill.background()

            # Label and percentage
            self.add_text_box(
                slide, 6.2, y_pos, 3.0, 0.5,
                f"{label}: {pct}%",
                font_size=14, bold=True, color='text'
            )

            # Count
            self.add_text_box(
                slide, 9.5, y_pos, 2.0, 0.5,
                f"({count} questions)",
                font_size=11, color='neutral'
            )

            y_pos += 0.7

        # Per-question sentiment (if available)
        if data['question_sentiments']:
            self.add_text_box(
                slide, 5.8, 4.5, 7.0, 0.4,
                "Per-Question Sentiment:",
                font_size=12, bold=True, color='secondary'
            )

            sentiment_bars_path = charts.sentiment_bars(
                data['question_sentiments'][:8],
                figsize=(5.5, 3)
            )
            self.add_chart_image(slide, sentiment_bars_path, 5.8, 4.9, 6.5, 2.3)

        charts.cleanup()
        return prs

    def generate_pdf_content(self, pdf: FPDF) -> FPDF:
        data = self._get_sentiment_data()

        self.pdf_add_content_page(pdf, "Sentiment Analysis")

        # Summary header
        pdf.set_font('Helvetica', 'B', 14)
        pdf.set_text_color(*self.pdf_colors['primary'])
        pdf.cell(0, 10, f"Analysis of {data['question_count']} Questions", new_x='LMARGIN', new_y='NEXT')
        pdf.ln(5)

        sentiments = [
            ('Positive', data['positive_pct'], data['positive_count'], self.pdf_colors['positive']),
            ('Neutral', data['neutral_pct'], data['neutral_count'], self.pdf_colors['neutral']),
            ('Negative', data['negative_pct'], data['negative_count'], self.pdf_colors['negative']),
        ]

        for label, pct, count, color in sentiments:
            # Color bar
            pdf.set_fill_color(*color)
            bar_width = pct * 1.5  # Scale percentage to width
            pdf.cell(bar_width, 12, '', fill=True)

            # Label
            pdf.set_font('Helvetica', 'B', 12)
            pdf.set_text_color(*self.pdf_colors['text'])
            pdf.cell(50, 12, f"  {label}: {pct}%")

            # Count
            pdf.set_font('Helvetica', '', 10)
            pdf.set_text_color(*self.pdf_colors['neutral'])
            pdf.cell(0, 12, f"({count} questions)", new_x='LMARGIN', new_y='NEXT')
            pdf.ln(2)

        # Per-question breakdown
        if data['question_sentiments']:
            pdf.ln(8)
            pdf.set_font('Helvetica', 'B', 11)
            pdf.set_text_color(*self.pdf_colors['secondary'])
            pdf.cell(0, 8, "Per-Question Sentiment:", new_x='LMARGIN', new_y='NEXT')

            sentiment_symbols = {
                'Positive': '+',
                'Neutral': '=',
                'Negative': '-',
            }

            for q_label, sentiment in data['question_sentiments'][:10]:
                symbol = sentiment_symbols.get(sentiment, '?')
                color = self.pdf_colors.get(sentiment.lower(), self.pdf_colors['neutral'])

                pdf.set_font('Helvetica', '', 9)
                pdf.set_text_color(*self.pdf_colors['text'])
                pdf.cell(120, 5, f"  {q_label}")

                pdf.set_font('Helvetica', 'B', 9)
                pdf.set_text_color(*color)
                pdf.cell(0, 5, f"[{symbol}] {sentiment}", new_x='LMARGIN', new_y='NEXT')

        return pdf
